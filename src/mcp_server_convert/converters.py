"""Document conversion engines."""

import csv
import io
import json
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any


def get_supported_formats() -> dict[str, dict[str, str]]:
    """Return dict of extension -> {description, method}."""
    formats = {
        "pdf": {"description": "PDF documents", "method": "PyMuPDF (fitz)"},
        "docx": {"description": "Microsoft Word documents", "method": "python-docx"},
        "html": {"description": "HTML web pages", "method": "markdownify"},
        "htm": {"description": "HTML web pages", "method": "markdownify"},
        "epub": {"description": "EPUB ebooks", "method": "ebooklib"},
        "csv": {"description": "CSV spreadsheets", "method": "stdlib csv"},
        "json": {"description": "JSON data", "method": "stdlib json"},
        "xml": {"description": "XML data", "method": "stdlib xml"},
        "xlsx": {"description": "Microsoft Excel spreadsheets", "method": "openpyxl"},
        "pptx": {"description": "Microsoft PowerPoint presentations", "method": "python-pptx"},
        "txt": {"description": "Plain text", "method": "passthrough"},
        "md": {"description": "Markdown", "method": "passthrough"},
        "rst": {"description": "reStructuredText", "method": "passthrough"},
        "log": {"description": "Log files", "method": "passthrough"},
        "yaml": {"description": "YAML data", "method": "formatted text"},
        "yml": {"description": "YAML data", "method": "formatted text"},
        "toml": {"description": "TOML data", "method": "formatted text"},
    }
    return formats


def sniff_format(path: Path) -> str:
    """Detect file format from extension."""
    ext = path.suffix.lower().lstrip(".")
    formats = get_supported_formats()
    return formats.get(ext, {}).get("description", f"Unknown (.{ext})")


def convert_file_to_markdown(path: Path) -> str:
    """Convert a file to Markdown based on its extension."""
    ext = path.suffix.lower().lstrip(".")

    converters = {
        "pdf": _convert_pdf,
        "docx": _convert_docx,
        "html": _convert_html,
        "htm": _convert_html,
        "epub": _convert_epub,
        "csv": _convert_csv,
        "json": _convert_json,
        "xml": _convert_xml,
        "xlsx": _convert_xlsx,
        "pptx": _convert_pptx,
        "yaml": _convert_text,
        "yml": _convert_text,
        "toml": _convert_text,
    }

    # Text passthrough
    text_exts = {"txt", "md", "rst", "log"}
    if ext in text_exts:
        return path.read_text(encoding="utf-8", errors="replace")

    converter = converters.get(ext)
    if converter:
        return converter(path)

    # Fallback: try reading as text
    try:
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception:
        raise ValueError(f"Unsupported file format: .{ext}")


def convert_bytes(data: bytes, filename: str) -> str:
    """Convert bytes to Markdown given a filename (for extension detection)."""
    ext = Path(filename).suffix.lower().lstrip(".")
    if ext in {"txt", "md", "rst", "log"}:
        return data.decode("utf-8", errors="replace")
    if ext == "json":
        obj = json.loads(data)
        return f"```json\n{json.dumps(obj, indent=2, ensure_ascii=False)}\n```"
    if ext == "csv":
        return _csv_to_md(data.decode("utf-8", errors="replace"))
    # For binary formats, save to temp and convert
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as f:
        f.write(data)
        f.flush()
        try:
            return convert_file_to_markdown(Path(f.name))
        finally:
            Path(f.name).unlink(missing_ok=True)


# --- Individual converters ---

def _convert_pdf(path: Path) -> str:
    """Convert PDF to Markdown using PyMuPDF."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return _convert_pdf_fallback(path)

    doc = fitz.open(str(path))
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            pages.append(f"## Page {i + 1}\n\n{text}")
    doc.close()
    return "\n\n---\n\n".join(pages) if pages else "(Empty PDF)"


def _convert_pdf_fallback(path: Path) -> str:
    """Fallback PDF extraction using basic text extraction."""
    try:
        # Try pdfminer
        from pdfminer.high_level import extract_text
        text = extract_text(str(path))
        return text if text.strip() else "(Empty or image-only PDF)"
    except ImportError:
        return "(PDF conversion requires PyMuPDF: pip install PyMuPDF)"


def _convert_docx(path: Path) -> str:
    """Convert DOCX to Markdown."""
    try:
        from docx import Document
    except ImportError:
        return "(DOCX conversion requires python-docx: pip install python-docx)"

    doc = Document(str(path))
    lines = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue

        style = para.style.name if para.style else ""

        if "Heading 1" in style:
            lines.append(f"# {text}")
        elif "Heading 2" in style:
            lines.append(f"## {text}")
        elif "Heading 3" in style:
            lines.append(f"### {text}")
        elif "Heading 4" in style:
            lines.append(f"#### {text}")
        elif "List" in style:
            lines.append(f"- {text}")
        else:
            lines.append(text)

    # Extract tables
    for table in doc.tables:
        lines.append("")
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("| " + " | ".join(["---"] * len(cells)) + " |")

    return "\n".join(lines)


def _convert_html(path: Path) -> str:
    """Convert HTML to Markdown."""
    try:
        from markdownify import markdownify
    except ImportError:
        return "(HTML conversion requires markdownify: pip install markdownify)"

    html = path.read_text(encoding="utf-8", errors="replace")
    return markdownify(html, heading_style="ATX")


def _convert_epub(path: Path) -> str:
    """Convert EPUB to Markdown."""
    try:
        import ebooklib
        from ebooklib import epub
        from markdownify import markdownify
    except ImportError:
        return "(EPUB conversion requires ebooklib and markdownify: pip install ebooklib markdownify)"

    book = epub.read_epub(str(path))
    chapters = []

    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        content = item.get_content().decode("utf-8", errors="replace")
        md = markdownify(content, heading_style="ATX").strip()
        if md:
            chapters.append(md)

    return "\n\n---\n\n".join(chapters)


def _convert_csv(path: Path) -> str:
    """Convert CSV to Markdown table."""
    text = path.read_text(encoding="utf-8", errors="replace")
    return _csv_to_md(text)


def _csv_to_md(text: str) -> str:
    """Convert CSV text to Markdown table."""
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)

    if not rows:
        return "(Empty CSV)"

    lines = []
    # Header
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
    # Data rows (limit to 100)
    for row in rows[1:101]:
        # Pad or truncate to match header length
        padded = row + [""] * (len(rows[0]) - len(row))
        lines.append("| " + " | ".join(padded[: len(rows[0])]) + " |")

    if len(rows) > 101:
        lines.append(f"\n... ({len(rows) - 101} more rows)")

    return "\n".join(lines)


def _convert_json(path: Path) -> str:
    """Convert JSON to formatted Markdown code block."""
    text = path.read_text(encoding="utf-8", errors="replace")
    obj = json.loads(text)
    formatted = json.dumps(obj, indent=2, ensure_ascii=False)
    return f"```json\n{formatted}\n```"


def _convert_xml(path: Path) -> str:
    """Convert XML to readable Markdown."""
    try:
        import xmltodict
        data = path.read_text(encoding="utf-8", errors="replace")
        d = xmltodict.parse(data)
        return f"```json\n{json.dumps(d, indent=2, ensure_ascii=False)}\n```"
    except ImportError:
        # Fallback: pretty print XML
        tree = ET.parse(path)
        root = tree.getroot()
        lines = [f"# {root.tag}", ""]

        def _walk(elem: ET.Element, depth: int = 0):
            indent = "  " * depth
            text = elem.text.strip() if elem.text else ""
            attrs = f" ({dict(elem.attrib)})" if elem.attrib else ""
            if text:
                lines.append(f"{indent}- **{elem.tag}**{attrs}: {text}")
            else:
                lines.append(f"{indent}- **{elem.tag}**{attrs}")
            for child in elem:
                _walk(child, depth + 1)

        _walk(root)
        return "\n".join(lines)


def _convert_xlsx(path: Path) -> str:
    """Convert XLSX to Markdown tables."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "(XLSX conversion requires openpyxl: pip install openpyxl)"

    wb = load_workbook(str(path), read_only=True)
    sheets = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        lines = [f"## Sheet: {sheet_name}", ""]

        # Header
        header = [str(c) if c is not None else "" for c in rows[0]]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")

        # Data (limit to 100 rows)
        for row in rows[1:101]:
            cells = [str(c) if c is not None else "" for c in row]
            padded = cells + [""] * (len(header) - len(cells))
            lines.append("| " + " | ".join(padded[: len(header)]) + " |")

        if len(rows) > 101:
            lines.append(f"\n... ({len(rows) - 101} more rows)")

        sheets.append("\n".join(lines))

    wb.close()
    return "\n\n---\n\n".join(sheets) if sheets else "(Empty spreadsheet)"


def _convert_pptx(path: Path) -> str:
    """Convert PPTX to Markdown slides."""
    try:
        from pptx import Presentation
    except ImportError:
        return "(PPTX conversion requires python-pptx: pip install python-pptx)"

    prs = Presentation(str(path))
    slides = []

    for i, slide in enumerate(prs.slides):
        lines = [f"## Slide {i + 1}", ""]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)

            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        lines.append("| " + " | ".join(["---"] * len(cells)) + " |")

        slides.append("\n".join(lines))

    return "\n\n---\n\n".join(slides)


def _convert_text(path: Path) -> str:
    """Passthrough for text files."""
    text = path.read_text(encoding="utf-8", errors="replace")
    ext = path.suffix.lower().lstrip(".")
    return f"```{ext}\n{text}\n```"
